from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')

def generar_reporte_pptx(datos, img_producto):
    """Función antigua - mantener para compatibilidad"""
    return generar_reporte_pptx_mejorado(datos, img_producto)

def generar_reporte_pptx_mejorado(datos, img_producto):
    """
    Genera reporte PowerPoint mejorado con:
    - Múltiples diapositivas
    - Gráficos
    - Análisis de capacidad
    - Branding corporativo
    """
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ====================================================================
    # DIAPOSITIVA 1: PORTADA
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # Azul oscuro
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = f"OFERTA FACTORY 21"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_text = f"Proyecto: {datos['proyecto']}\nOEM: {datos['oem']} | Versión: {datos['version']}"
    subtitle_frame.text = subtitle_text
    for para in subtitle_frame.paragraphs:
        para.font.size = Pt(24)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
    
    # ====================================================================
    # DIAPOSITIVA 2: RESUMEN EJECUTIVO
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Resumen Ejecutivo"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    # Agregar contenido
    p = tf.paragraphs[0]
    p.text = f"Tiempo de Ciclo: {datos['t_ciclo']:.2f} segundos"
    p.font.size = Pt(18)
    p.font.bold = True
    
    items = [
        f"MOD por Instalación: {datos['n_mod']} módulo(s)",
        f"Saturación Manual: {datos['saturacion']:.1f}%",
        f"Capacidad Máxima: {datos.get('cap_max', 0):,.0f} piezas/año",
        f"Modelo: Regresión Lineal Múltiple",
        f"Factor IA: {datos['factor_ia']:.3f}"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.level = 1
    
    # Insertar imagen si existe
    if img_producto:
        try:
            slide.shapes.add_picture(BytesIO(img_producto.read()), Inches(5.5), Inches(2), width=Inches(4))
        except:
            pass
    
    # ====================================================================
    # DIAPOSITIVA 3: ANÁLISIS TÉCNICO
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Análisis Técnico - Tiempo de Ciclo"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = f"Composición del Tiempo de Ciclo"
    p.font.size = Pt(16)
    p.font.bold = True
    
    # Cálculo de componentes
    t_ciclo = datos['t_ciclo']
    t_sold = datos.get('t_soldadores', t_ciclo * 0.7)
    t_manip = datos.get('t_manipulador', t_ciclo * 0.3)
    
    pct_sold = (t_sold / t_ciclo * 100) if t_ciclo > 0 else 0
    pct_manip = (t_manip / t_ciclo * 100) if t_ciclo > 0 else 0
    
    items = [
        f"Soldadores: {t_sold:.2f}s ({pct_sold:.1f}%)",
        f"Manipulador: {t_manip:.2f}s ({pct_manip:.1f}%)",
        f"Total Ciclo: {t_ciclo:.2f}s"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.level = 1
    
    # Crear gráfico circular
    try:
        fig, ax = plt.subplots(figsize=(4, 3))
        tiempos = [t_sold, t_manip]
        labels = ['Soldadores', 'Manipulador']
        colors = ['#FF6B6B', '#4ECDC4']
        ax.pie(tiempos, labels=labels, autopct='%1.1f%%', colors=colors, startangle=90)
        ax.set_title('Distribución de Tiempo')
        
        img_buffer = BytesIO()
        fig.savefig(img_buffer, format='png', bbox_inches='tight', dpi=100)
        img_buffer.seek(0)
        
        slide.shapes.add_picture(img_buffer, Inches(5.5), Inches(2), width=Inches(4))
        plt.close(fig)
    except Exception as e:
        print(f"Error creando gráfico: {e}")
    
    # ====================================================================
    # DIAPOSITIVA 4: ANÁLISIS DE CAPACIDAD
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Plan de Capacidad"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Necesidades de Instalaciones por Año"
    p.font.size = Pt(16)
    p.font.bold = True
    
    # Datos de capacidad
    if 'res_anual' in datos and datos['res_anual']:
        for item in datos['res_anual']:
            p = tf.add_paragraph()
            p.text = f"Año {item['Año']}: {item['Volumen']:,} piezas → {item['Instalaciones']} línea(s) ({item['Operarios_Turno']} operarios/turno)"
            p.font.size = Pt(12)
            p.level = 1
    
    # Recomendaciones
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Recomendaciones:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    saturacion = datos['saturacion']
    if saturacion < 50:
        recom = "✓ Instalar 1 MOD (operario dedicado)"
    elif saturacion < 100:
        recom = "✓ Instalar 1 MOD con buffer de capacidad"
    else:
        recom = f"✓ Instalar {datos['n_mod']} MOD (saturación > 100%)"
    
    p = tf.add_paragraph()
    p.text = recom
    p.font.size = Pt(12)
    p.level = 1
    
    # ====================================================================
    # DIAPOSITIVA 5: INFORMACIÓN DEL MODELO
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Modelo Matemático"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Técnica: Regresión Lineal Múltiple"
    p.font.size = Pt(14)
    p.font.bold = True
    
    items = [
        "Algoritmo: Mínimos Cuadrados Ordinarios (OLS)",
        "Datos de Entrenamiento: 8 proyectos históricos",
        "Variables Independientes: SPW, Peso, ANCHO_ASSY",
        f"R² Score: {datos.get('r2_score', 0.70):.2%} (varianza explicada)",
        f"RMSE: ~14.80 segundos (error típico)"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.level = 1
    
    # ====================================================================
    # DIAPOSITIVA 6: NOTAS IMPORTANTES
    # ====================================================================
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Notas Importantes"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "⚠️ Consideraciones y Limitaciones"
    p.font.size = Pt(14)
    p.font.bold = True
    
    items = [
        "El modelo se basa en datos históricos de procesos similares.",
        "La precisión depende de la consistencia de datos y procesos.",
        "Se recomienda validar con pruebas piloto antes de inversión.",
        f"Margen de error estimado: ±{datos.get('mae', 11):.1f} segundos",
        "Para procesamiento manual (MOD), aplicar factores de seguridad."
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.level = 1
    
    # ====================================================================
    # GUARDAR PRESENTACIÓN
    # ====================================================================
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.getvalue()
